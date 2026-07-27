"""
ic_deck_generator.py — Investment Committee deck (D30/D43; rebuilt for the
9-section professional IC format — Tâche "Mémo IC et deck au format IC
professionnel").

Mirrors the 9-section Word memo (`docx_generator.py`) slide by slide, both
consuming the SAME structured context (`ic_context.build_ic_context`) so the
memo and the deck can never disagree on a figure, a recommendation, or a
risk. Design system (D43) — light/sober palette, semantic green/red
reserved for IRR/MOIC — is unchanged; only slide content/structure changes.

Point d'entrée :
    generate_ic_deck(deal, comps_table=None, reference_scenario=None, ic_context=None) -> io.BytesIO
"""
from __future__ import annotations

import io
import re
from datetime import datetime

from loguru import logger
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from api.services.ma_engine.ic_context import (
    MEMO_SECTIONS,
    build_ic_context,
    extract_markdown_section,
    strip_markdown_bold,
)

# ============================================================
# Design system (D43) — light palette, 3 colours + neutrals.
# Green/red are SEMANTIC colours (IRR/MOIC performance only), not a
# decorative 4th/5th palette colour.
# ============================================================

COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BG_ALT = RGBColor(0xF6, 0xF7, 0xF9)
COLOR_INK = RGBColor(0x1A, 0x22, 0x33)
COLOR_MUTED = RGBColor(0x6B, 0x72, 0x80)
COLOR_ACCENT = RGBColor(0x1F, 0x4E, 0x79)
COLOR_ACCENT_LIGHT = RGBColor(0xE8, 0xEF, 0xF6)
COLOR_LINE = RGBColor(0xDD, 0xE1, 0xE7)
COLOR_GOOD = RGBColor(0x1F, 0x7A, 0x4D)
COLOR_BAD = RGBColor(0xB1, 0x3B, 0x3B)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Arial"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MARGIN_L = 0.7
CONTENT_W = 11.93
HEADER_TOP = 0.55
HEADER_RULE_TOP = 1.35
CONTENT_TOP = 1.65
FOOTER_TOP = 7.05


# ============================================================
# Generic helpers
# ============================================================

def _no_line(shape) -> None:
    shape.line.fill.background()


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG
    return slide


def _text(slide, left, top, width, height, text, *, size=14, color=COLOR_INK,
          bold=False, italic=False, align=PP_ALIGN.LEFT, font=FONT, wrap=True,
          anchor=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    return box


def _multiline(slide, left, top, width, height, lines, *, size=12, color=COLOR_INK,
               line_spacing=1.2, bullet_color=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
    return box


def _rect(slide, left, top, width, height, *, fill=COLOR_BG_ALT, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = COLOR_LINE
        shape.line.width = Pt(0.75)
    else:
        _no_line(shape)
    shape.shadow.inherit = False
    return shape


def _hairline(slide, left, top, width, *, color=COLOR_LINE, weight=0.75):
    shape = slide.shapes.add_connector(1, Inches(left), Inches(top), Inches(left + width), Inches(top))
    shape.line.color.rgb = color
    shape.line.width = Pt(weight)
    return shape


def _kpi(slide, left, top, width, label, value, *, value_color=COLOR_ACCENT, sub=None,
         height=1.25, value_size=22):
    _rect(slide, left, top, width, height, fill=COLOR_BG_ALT)
    _hairline(slide, left, top, width, color=COLOR_ACCENT, weight=1.75)
    _text(slide, left + 0.18, top + 0.16, width - 0.36, 0.28, label.upper(),
          size=9, color=COLOR_MUTED, bold=True)
    value_zone_top = top + 0.5
    value_zone_bottom = top + height - (0.4 if sub else 0.12)
    _text(slide, left + 0.18, value_zone_top, width - 0.36, value_zone_bottom - value_zone_top, value,
          size=value_size, color=value_color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        _text(slide, left + 0.18, top + height - 0.32, width - 0.36, 0.28, sub,
              size=8, color=COLOR_MUTED, italic=True)


def _header(slide, kicker, title, subtitle=None):
    _text(slide, MARGIN_L, HEADER_TOP - 0.32, CONTENT_W, 0.26, kicker.upper(),
          size=9, color=COLOR_ACCENT, bold=True)
    # Tâche "P3" (Partie A/D) — les titres d'action sont des phrases
    # calculées (longueur variable selon le nom de la cible/du secteur),
    # contrairement à l'ancien intitulé de section fixe. Un titre trop long
    # à taille fixe passait à la ligne et chevauchait le sous-titre/la règle
    # (constaté à la vérification visuelle) — la taille est réduite en
    # fonction de la longueur plutôt que de laisser un débordement silencieux.
    if len(title) > 82:
        title_size = 17
    elif len(title) > 68:
        title_size = 19
    elif len(title) > 55:
        title_size = 21
    else:
        title_size = 23
    _text(slide, MARGIN_L, HEADER_TOP, CONTENT_W, 0.5, title, size=title_size, color=COLOR_INK, bold=True)
    if subtitle:
        _text(slide, MARGIN_L, HEADER_TOP + 0.46, CONTENT_W, 0.3, subtitle,
              size=11, color=COLOR_MUTED, italic=True)
    _hairline(slide, MARGIN_L, HEADER_RULE_TOP, CONTENT_W, color=COLOR_ACCENT, weight=2.0)


def _footer(slide, page_no, page_total):
    _hairline(slide, MARGIN_L, FOOTER_TOP, CONTENT_W, color=COLOR_LINE, weight=0.75)
    _text(slide, MARGIN_L, FOOTER_TOP + 0.06, 6, 0.3, "PE Tracker — Confidential",
          size=8, color=COLOR_MUTED)
    _text(slide, MARGIN_L + CONTENT_W - 2, FOOTER_TOP + 0.06, 2, 0.3, f"{page_no} / {page_total}",
          size=8, color=COLOR_MUTED, align=PP_ALIGN.RIGHT)


def _source_note(slide, left, top, width, text):
    _text(slide, left, top, width, 0.3, text, size=8, color=COLOR_MUTED, italic=True)


def _money(value: float | None) -> str:
    if value is None:
        return "N/A"
    return f"€{value:,.0f}".replace(",", " ")


def _pct(value: float | None, *, already_pct: bool = False) -> str:
    if value is None:
        return "N/A"
    v = value if already_pct else value * 100
    return f"{v:.1f}%"


_LIST_ITEM_RE = re.compile(r"^\s*(?:[-•*]|\d+[.)])\s+")


def _bullets_from_section(markdown_text: str | None, heading: str, *, marker: str = "•", limit: int = 10) -> list[str]:
    """Extracts a section's lines from the memo and formats them as a
    uniform bullet list (D43 : one bullet style per slide, never a mix of
    the memo's own '-'/plain-text/numbered-list conventions)."""
    lines = extract_markdown_section(markdown_text, heading)
    if not lines:
        return [f"IC memo narrative not available — generate the memo first."]
    out = []
    for line in lines[:limit]:
        stripped = strip_markdown_bold(line).lstrip()
        if stripped.startswith(("-", "•", "⚠")):
            stripped = stripped[1:].lstrip()
        stripped = re.sub(r"^\d+[.)]\s+", "", stripped)
        out.append(f"{marker} {stripped}")
    return out


def _marked_bullets_from_section(markdown_text: str | None, heading: str, *, marker: str = "⚠", limit: int = 10) -> list[str]:
    """Tâche "P3 : le deck IC devient un vrai deck" (Partie D) — même
    extraction que `_bullets_from_section`, mais le marqueur n'est appliqué
    QU'AUX lignes qui sont réellement des puces dans le markdown source
    (préfixées par -/•/*/1.). Une ligne d'introduction en prose libre (ex.
    « Key risks identified include: ») n'a jamais été une puce et ne doit
    jamais porter un ⚠ — avant ce correctif, `_bullets_from_section`
    préfixait INDISCRIMINÉMENT chaque ligne renvoyée par `marker`, y compris
    une éventuelle ligne d'intro non listée (constat de la revue IC externe)."""
    lines = extract_markdown_section(markdown_text, heading)
    if not lines:
        return [f"IC memo narrative not available — generate the memo first."]
    out = []
    for line in lines[:limit]:
        is_list_item = bool(_LIST_ITEM_RE.match(line))
        stripped = strip_markdown_bold(line).lstrip()
        if stripped.startswith(("-", "•", "⚠", "*")):
            stripped = stripped[1:].lstrip()
        stripped = re.sub(r"^\d+[.)]\s+", "", stripped)
        out.append(f"{marker} {stripped}" if is_list_item else stripped)
    return out


def _condensed_fragments_from_section(markdown_text: str | None, heading: str, *,
                                       marker: str = "•", max_source_lines: int = 2, max_fragments: int = 6) -> list[str]:
    """Tâche "P3 : le deck IC devient un vrai deck" (Partie C) — condense un
    paragraphe de prose du mémo en fragments de phrase courts au lieu de
    coller le paragraphe brut en une seule puce géante. Constat de la revue
    IC externe : les slides Executive Summary / Company Overview collaient
    la prose du mémo telle quelle, produisant une "slide à un seul
    paragraphe" — un deck se commente, il ne se lit pas."""
    prose_lines = [
        strip_markdown_bold(l) for l in extract_markdown_section(markdown_text, heading)
        if not l.lstrip().startswith(("-", "*", "•", "⚠"))
    ]
    if not prose_lines:
        return [f"{marker} IC memo narrative not available — generate the memo first."]
    fragments = []
    for line in prose_lines[:max_source_lines]:
        for frag in re.split(r"(?<=[.!?])\s+", line):
            frag = frag.strip()
            if frag:
                fragments.append(f"{marker} {frag}")
    return fragments[:max_fragments] or [f"{marker} IC memo narrative not available — generate the memo first."]


def _adaptive_bullets_from_section(markdown_text: str | None, heading: str, *,
                                    marker: str = "•", limit: int = 6, max_fragments: int = 6) -> list[str]:
    """Tâche "P3" (Partie C) — le LLM rédige certaines sections du mémo en
    liste à puces et d'autres en prose libre, et ce choix varie d'un deal à
    l'autre (ex. Section V "Investment Thesis" : bullets pour BTP, un seul
    paragraphe de prose pour Ingebime). Choisir la mauvaise stratégie de
    rendu produit soit une liste tronquée soit — le bug ici — une slide à
    un seul paragraphe. Ce dispatcher détecte la forme réelle du texte
    source et applique le rendu adapté plutôt que de supposer un format fixe."""
    lines = extract_markdown_section(markdown_text, heading)
    if not lines:
        return [f"IC memo narrative not available — generate the memo first."]
    list_item_count = sum(1 for l in lines if _LIST_ITEM_RE.match(l))
    if list_item_count >= max(1, len(lines) // 2):
        return _bullets_from_section(markdown_text, heading, marker=marker, limit=limit)
    return _condensed_fragments_from_section(markdown_text, heading, marker=marker, max_source_lines=3, max_fragments=max_fragments)


# ============================================================
# Table helper — contrasted header, alternating rows, no heavy grid.
# ============================================================

def _styled_table(slide, rows, cols, left, top, width, height, headers, data,
                   *, col_align=None, header_size=10.5, body_size=10, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        # Tâche "P3" (Partie D) — colonnes proportionnelles au contenu plutôt
        # qu'équi-réparties, pour laisser assez de place à une colonne texte
        # (ex. "Reference") sans avoir à tronquer la chaîne en plein mot.
        for j, w in enumerate(col_widths):
            table.columns[j].width = Inches(w)

    tbl = table_shape.table._tbl
    tbl_pr = tbl.find(qn('a:tblPr'))
    if tbl_pr is not None:
        tbl_pr.set('firstRow', '0')
        tbl_pr.set('bandRow', '0')

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(header_size)
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT
        if col_align and col_align[j] == "right":
            p.alignment = PP_ALIGN.RIGHT
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_ACCENT
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.12)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, row_vals in enumerate(data, start=1):
        row_fill = COLOR_BG if i % 2 else COLOR_BG_ALT
        for j, v in enumerate(row_vals):
            cell = table.cell(i, j)
            cell.text = str(v)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(body_size)
            p.font.color.rgb = COLOR_INK
            p.font.name = FONT
            if col_align and col_align[j] == "right":
                p.alignment = PP_ALIGN.RIGHT
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_fill
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


# ============================================================
# Native charts (Tâche "P3 : le deck IC devient un vrai deck", Partie B) —
# objets python-pptx natifs (graphique + feuille de données embarquée),
# jamais une image figée : chaque graphique reste éditable dans PowerPoint.
# Toutes les données viennent du MÊME scenario/context que le mémo (D23) —
# aucun nouveau calcul, uniquement une mise en forme visuelle des chiffres
# déjà produits par le moteur/la calibration.
# ============================================================

def _label_color_for_bg(rgb: RGBColor) -> RGBColor:
    """Tâche "P3" (Partie D) — un label de données dans la même couleur
    (sombre) que son propre fond de barre est invisible (constaté : la
    barre "Sponsor Equity"/"Exit Equity", remplie de COLOR_INK, portait un
    label COLOR_INK — illisible). Choix blanc/encre par luminance perçue,
    jamais une couleur de label fixe indépendante du fond réel de la barre."""
    r, g, b = rgb[0], rgb[1], rgb[2]
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return COLOR_WHITE if luminance < 140 else COLOR_INK


def _style_chart_base(chart, *, has_legend=False, legend_pos=XL_LEGEND_POSITION.BOTTOM):
    chart.has_legend = has_legend
    if has_legend:
        chart.legend.position = legend_pos
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = FONT
    chart.has_title = False
    try:
        chart.category_axis.tick_labels.font.size = Pt(9)
        chart.category_axis.tick_labels.font.name = FONT
        chart.category_axis.format.line.color.rgb = COLOR_LINE
        chart.category_axis.major_tick_mark = 0
    except Exception:
        pass
    try:
        chart.value_axis.tick_labels.font.size = Pt(9)
        chart.value_axis.tick_labels.font.name = FONT
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = COLOR_LINE
        chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
        chart.value_axis.format.line.fill.background()
    except Exception:
        pass


def _add_column_chart(slide, left, top, width, height, categories, series, *,
                       colors=None, number_format='#,##0', stacked=False, has_legend=None,
                       data_labels=True, label_number_format=None, value_axis_min=None):
    """Clustered or stacked column chart. `series` = [(name, [values...]), ...] —
    or [(name, [values...], show_labels_bool), ...] to suppress data labels on
    a specific series (Tâche "P3", Partie D : sur un graphique Sources & Uses
    empilé, une petite tranche de frais affiche un label qui chevauche celui
    du segment voisin — désactiver le label des tranches non essentielles au
    lieu de laisser un chevauchement, constaté à la vérification visuelle)."""
    cd = CategoryChartData()
    cd.categories = categories
    show_labels_per_series = []
    for item in series:
        name, values = item[0], item[1]
        show_labels_per_series.append(item[2] if len(item) > 2 else True)
        cd.add_series(name, values)
    chart_type = XL_CHART_TYPE.COLUMN_STACKED if stacked else XL_CHART_TYPE.COLUMN_CLUSTERED
    gframe = slide.shapes.add_chart(chart_type, Inches(left), Inches(top), Inches(width), Inches(height), cd)
    chart = gframe.chart
    if has_legend is None:
        has_legend = len(series) > 1
    _style_chart_base(chart, has_legend=has_legend)
    plot = chart.plots[0]
    plot.gap_width = 60 if not stacked else 90
    if len(series) > 1 and not stacked:
        plot.overlap = -10
    default_colors = colors or [COLOR_ACCENT, COLOR_MUTED, COLOR_GOOD, COLOR_BAD, COLOR_INK]
    for i, s in enumerate(plot.series):
        series_color = default_colors[i % len(default_colors)]
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = series_color
        s.format.line.fill.background()
        if data_labels and show_labels_per_series[i]:
            s.data_labels.show_value = True

            # Tâche "P3" (Partie D) — un format à 3 sections (positif;négatif;
            # zéro="") masque les labels "0" sur les points où cette série ne
            # contribue pas à la catégorie (ex. "Enterprise Value" vaut 0 côté
            # "Sources" sur un graphique Sources & Uses empilé) — sans cela,
            # un "0" flottant apparaissait en bas de chaque colonne empilée
            # (constaté à la vérification visuelle), pas un vrai chiffre utile.
            base_fmt = label_number_format or number_format
            s.data_labels.number_format = f'{base_fmt};-{base_fmt};""'
            s.data_labels.number_format_is_linked = False
            s.data_labels.font.size = Pt(9)
            s.data_labels.font.bold = True
            s.data_labels.font.name = FONT
            # Tâche "P3" (Partie D) — la luminance ne compte QUE sur un
            # graphique empilé : le label y est centré SUR le remplissage de
            # la barre (un label sombre sur une barre sombre, ex. "Sponsor
            # Equity" en COLOR_INK, était invisible). Sur un graphique en
            # colonnes groupées (non empilé), le label par défaut se pose
            # AU-DESSUS de la barre, sur le fond blanc de la slide — y
            # appliquer la même règle de luminance rendait le label blanc
            # invisible sur fond blanc (constaté : les barres Revenue/EBITDA
            # de Financial Analysis avaient perdu leurs labels).
            s.data_labels.font.color.rgb = _label_color_for_bg(series_color) if stacked else COLOR_INK
    try:
        chart.value_axis.tick_labels.number_format = number_format
        chart.value_axis.tick_labels.number_format_is_linked = False
        # Tâche "P3" (Partie D) — un axe qui ne part pas de 0 exagère
        # visuellement un écart entre deux barres (constaté sur le
        # graphique Base vs Downside MOIC, où l'auto-échelle démarrait à
        # 1.5x) — jamais laissé au hasard de l'auto-échelle pour une
        # comparaison où l'honnêteté visuelle de l'écart compte.
        if value_axis_min is not None:
            chart.value_axis.minimum_scale = value_axis_min
    except Exception:
        pass
    return chart


def _add_line_chart(slide, left, top, width, height, categories, series, *,
                     colors=None, number_format='0.0"x"'):
    """Line-with-markers chart — used for the deleveraging curve. `series` =
    [(name, [values...]), ...] or [(name, [values...], show_labels_bool), ...]
    — two near-identical lines (net vs gross leverage) both labelled produce
    illegible overlapping text (constaté à la vérification visuelle) ; the
    secondary line can opt out of data labels while keeping its own line/marker."""
    cd = CategoryChartData()
    cd.categories = categories
    show_labels_per_series = []
    for item in series:
        name, values = item[0], item[1]
        show_labels_per_series.append(item[2] if len(item) > 2 else True)
        cd.add_series(name, values)
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(left), Inches(top), Inches(width), Inches(height), cd,
    )
    chart = gframe.chart
    _style_chart_base(chart, has_legend=len(series) > 1)
    default_colors = colors or [COLOR_ACCENT, COLOR_MUTED]
    for i, s in enumerate(chart.plots[0].series):
        s.format.line.color.rgb = default_colors[i % len(default_colors)]
        s.format.line.width = Pt(2.25)
        s.marker.style = 8  # circle
        s.marker.format.fill.solid()
        s.marker.format.fill.fore_color.rgb = default_colors[i % len(default_colors)]
        s.marker.format.line.color.rgb = default_colors[i % len(default_colors)]
        if not show_labels_per_series[i]:
            continue
        s.data_labels.show_value = True

        s.data_labels.number_format = number_format
        s.data_labels.number_format_is_linked = False
        s.data_labels.font.size = Pt(9)
        s.data_labels.font.bold = True
        s.data_labels.font.name = FONT
        s.data_labels.font.color.rgb = default_colors[i % len(default_colors)]
    try:
        chart.value_axis.tick_labels.number_format = number_format
        chart.value_axis.tick_labels.number_format_is_linked = False
    except Exception:
        pass
    return chart


def _add_bridge_chart(slide, left, top, width, height, steps, *, number_format='#,##0',
                       start_color=COLOR_ACCENT, end_color=COLOR_INK,
                       up_color=COLOR_GOOD, down_color=COLOR_BAD):
    """Waterfall/bridge chart via the classic stacked-column trick (an
    invisible 'base' series carries the floating bar's floor, a visible
    'value' series draws the step) — `steps` = [(label, value, kind), ...]
    where kind ∈ {"total", "increase", "decrease"}. The FIRST and LAST steps
    must be "total" (full bars from 0); every step in between is a delta
    applied to the running total. Never a new number: `value` for a delta
    step is the actual change already implied by the underlying data."""
    categories = [s[0] for s in steps]
    base_vals: list[float] = []
    value_vals: list[float] = []
    running = 0.0
    for label, value, kind in steps:
        if kind == "total":
            base_vals.append(0.0)
            value_vals.append(value)
            running = value
        else:
            new_running = running + value if kind == "increase" else running - abs(value)
            base_vals.append(min(running, new_running))
            value_vals.append(abs(new_running - running))
            running = new_running

    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("base", base_vals)
    cd.add_series("value", value_vals)
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(left), Inches(top), Inches(width), Inches(height), cd,
    )
    chart = gframe.chart
    _style_chart_base(chart, has_legend=False)
    plot = chart.plots[0]
    plot.gap_width = 40
    base_series, value_series = plot.series[0], plot.series[1]
    base_series.format.fill.background()
    base_series.format.line.fill.background()
    value_series.format.line.fill.background()
    value_series.data_labels.show_value = True

    value_series.data_labels.number_format = number_format
    value_series.data_labels.number_format_is_linked = False
    value_series.data_labels.font.size = Pt(9.5)
    value_series.data_labels.font.bold = True
    value_series.data_labels.font.name = FONT
    value_series.data_labels.font.color.rgb = COLOR_INK
    for idx, (_, value, kind) in enumerate(steps):
        pt = value_series.points[idx]
        pt.format.fill.solid()
        if kind == "total":
            pt_color = start_color if idx == 0 else end_color
        elif kind == "increase":
            pt_color = up_color
        else:
            pt_color = down_color
        pt.format.fill.fore_color.rgb = pt_color
        # Tâche "P3" (Partie D) — un bâton sombre ("total"/COLOR_INK, ex.
        # "Exit Equity") portait un label dans la MÊME couleur sombre que son
        # propre fond — invisible à l'écran (constaté à la vérification
        # visuelle). Choix blanc/encre par luminance réelle du fond du point.
        pt.data_label.font.color.rgb = _label_color_for_bg(pt_color)
        pt.data_label.font.size = Pt(9.5)
        pt.data_label.font.bold = True
        pt.data_label.font.name = FONT
    try:
        chart.value_axis.tick_labels.number_format = number_format
        chart.value_axis.tick_labels.number_format_is_linked = False
    except Exception:
        pass
    return chart


def _add_range_bar_chart(slide, left, top, width, height, rows, *, number_format='0.0"x"',
                          base_color=COLOR_MUTED, highlight_color=COLOR_ACCENT):
    """Horizontal bar chart of one value per row (e.g. EV/EBITDA per listed
    comparable) — a "football-field-style" view (single point per company,
    not a per-method range, since that is the real data available), with the
    target's own retained multiple highlighted in a distinct colour.
    `rows` = [(label, value, is_target), ...]."""
    categories = [r[0] for r in rows]
    values = [r[1] for r in rows]
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("value", values)
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(left), Inches(top), Inches(width), Inches(height), cd,
    )
    chart = gframe.chart
    _style_chart_base(chart, has_legend=False)
    plot = chart.plots[0]
    plot.gap_width = 40
    series = plot.series[0]
    series.format.line.fill.background()
    series.data_labels.show_value = True

    series.data_labels.number_format = number_format
    series.data_labels.number_format_is_linked = False
    series.data_labels.font.size = Pt(9)
    series.data_labels.font.bold = True
    series.data_labels.font.name = FONT
    series.data_labels.font.color.rgb = COLOR_INK
    for idx, (_, _, is_target) in enumerate(rows):
        pt = series.points[idx]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = highlight_color if is_target else base_color
    try:
        chart.value_axis.tick_labels.number_format = number_format
        chart.value_axis.tick_labels.number_format_is_linked = False
    except Exception:
        pass
    return chart


# ============================================================
# Slides — one (or more) per IC memo section, same data context as the docx.
# ============================================================

def _slide_cover(prs: Presentation, deal) -> None:
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, 0.18, 7.5, fill=COLOR_ACCENT)
    _rect(slide, 0.18, 0, 0.04, 7.5, fill=COLOR_ACCENT_LIGHT)

    _text(slide, 0.9, 0.7, 11, 0.35, "INVESTMENT COMMITTEE", size=12, color=COLOR_ACCENT, bold=True)
    _hairline(slide, 0.9, 1.05, 3.2, color=COLOR_ACCENT, weight=2.0)

    _text(slide, 0.9, 2.7, 11.3, 1.4, deal.target_name or "Unnamed Target",
          size=42, color=COLOR_INK, bold=True)
    _text(slide, 0.9, 3.75, 11.3, 0.5,
          f"{deal.deal_type or 'M&A'} — {deal.sector or 'Sector not disclosed'}",
          size=15, color=COLOR_ACCENT, italic=True)

    _hairline(slide, 0.9, 4.5, 6.0, color=COLOR_LINE, weight=1.0)
    _text(slide, 0.9, 4.7, 11, 0.4,
          f"Sector: {deal.sector or 'N/A'}    |    Country: {deal.country or 'N/A'}    |    "
          f"{datetime.now().strftime('%d %B %Y')}",
          size=12, color=COLOR_MUTED)
    _text(slide, 0.9, 6.9, 6, 0.3, "PE Tracker — Confidential", size=9, color=COLOR_MUTED)


# Tâche "P3 : le deck IC devient un vrai deck" (Partie C) — slide Agenda,
# ajoutée en tête : un comité lit d'abord la structure de la présentation,
# jamais présente dans l'ancien deck (un simple empilement de sections).
def _slide_agenda(prs: Presentation, agenda_items: list[str], page, total) -> None:
    slide = _blank_slide(prs)
    _header(slide, "Investment Committee", "Agenda")
    row_h = 0.5
    top = CONTENT_TOP + 0.2
    for i, item in enumerate(agenda_items):
        y = top + i * row_h
        _text(slide, MARGIN_L, y, 0.5, row_h, f"{i + 1:02d}", size=15, color=COLOR_ACCENT, bold=True)
        _text(slide, MARGIN_L + 0.6, y, CONTENT_W - 0.6, row_h, item, size=15, color=COLOR_INK)
        if i < len(agenda_items) - 1:
            _hairline(slide, MARGIN_L, y + row_h - 0.06, CONTENT_W, color=COLOR_LINE, weight=0.5)
    _footer(slide, page, total)


def _slide_executive_summary(prs: Presentation, deal, ic_context, action_title, page, total) -> None:
    slide = _blank_slide(prs)
    _header(slide, "IC Memo — Section I", action_title)

    fin = ic_context["financials"]
    scenario = ic_context.get("reference_lbo_scenario")
    kpi_w = (CONTENT_W - 0.6) / 4
    row1 = CONTENT_TOP
    _kpi(slide, MARGIN_L, row1, kpi_w, "Revenue", _money(deal.target_revenue), value_color=COLOR_ACCENT, height=1.15, value_size=22)
    _kpi(slide, MARGIN_L + (kpi_w + 0.2), row1, kpi_w, "EBITDA", _money(deal.target_ebitda), value_color=COLOR_ACCENT, height=1.15, value_size=22)
    irr = scenario.get("irr") if scenario else None
    moic = scenario.get("moic") if scenario else None
    irr_color = COLOR_GOOD if irr and irr >= 0.20 else COLOR_BAD if irr is not None else COLOR_MUTED
    moic_color = COLOR_GOOD if moic and moic >= 2.5 else COLOR_BAD if moic is not None else COLOR_MUTED
    _kpi(slide, MARGIN_L + 2 * (kpi_w + 0.2), row1, kpi_w, "Projected IRR", _pct(irr) if irr is not None else "N/A", value_color=irr_color, height=1.15, value_size=22)
    _kpi(slide, MARGIN_L + 3 * (kpi_w + 0.2), row1, kpi_w, "Projected MOIC", f"{moic:.2f}x" if moic is not None else "N/A", value_color=moic_color, height=1.15, value_size=22)

    lines = _adaptive_bullets_from_section(deal.ic_memo, MEMO_SECTIONS[0], limit=6, max_fragments=6)
    _multiline(slide, MARGIN_L, row1 + 1.4, CONTENT_W, 4.0, lines, size=13.5)
    _footer(slide, page, total)


def _slide_company_overview(prs: Presentation, deal, ic_context, action_title, page, total) -> None:
    slide = _blank_slide(prs)
    _header(slide, "IC Memo — Section II", action_title)
    co = ic_context["company_overview"]

    # Renders the LLM's own English narrative for this section (same text as
    # the docx), NOT the raw `deal.description` — that field may be in the
    # deal's original language (French sourcing data) and would otherwise
    # make the deck disagree with the memo's own Section II wording and
    # break the English-only rule (D17) for this slide only. Condensed to
    # short fragments (Tâche "P3", Partie C) rather than pasted flowing
    # prose — a deck slide is commented on, not read.
    _text(slide, MARGIN_L, CONTENT_TOP, 4, 0.3, "BUSINESS DESCRIPTION", size=10, color=COLOR_MUTED, bold=True)
    condensed = _adaptive_bullets_from_section(deal.ic_memo, MEMO_SECTIONS[1], limit=4, max_fragments=4)
    _multiline(slide, MARGIN_L, CONTENT_TOP + 0.35, CONTENT_W, 1.7, condensed, size=13)

    kpi_top = 3.85
    kpi_w = (CONTENT_W - 0.4) / 3
    _kpi(slide, MARGIN_L, kpi_top, kpi_w, "Sector", co["sector"], value_color=COLOR_INK, value_size=17)
    _kpi(slide, MARGIN_L + kpi_w + 0.2, kpi_top, kpi_w, "Country", co["country"], value_color=COLOR_INK)
    _kpi(slide, MARGIN_L + 2 * (kpi_w + 0.2), kpi_top, kpi_w, "Deal Type", deal.deal_type or "N/A", value_color=COLOR_INK)

    notes_top = kpi_top + 1.45
    _text(slide, MARGIN_L, notes_top, 4, 0.3, "MANAGEMENT TEAM", size=10, color=COLOR_MUTED, bold=True)
    _multiline(slide, MARGIN_L, notes_top + 0.32, CONTENT_W / 2 - 0.2, 0.8, [co["management_team"]], size=11, color=COLOR_MUTED)
    _text(slide, MARGIN_L + CONTENT_W / 2 + 0.2, notes_top, 4, 0.3, "HEADCOUNT", size=10, color=COLOR_MUTED, bold=True)
    _multiline(slide, MARGIN_L + CONTENT_W / 2 + 0.2, notes_top + 0.32, CONTENT_W / 2 - 0.2, 0.8, [co["headcount"]], size=11, color=COLOR_MUTED)
    _footer(slide, page, total)


def _slide_industry_market_thesis(prs: Presentation, deal, ic_context, action_title, page, total) -> None:
    """Tâche "P3" (Partie C) — Section III (Industry & Market) et Section V
    (Investment Thesis) FUSIONNÉES sur une seule slide, en deux colonnes :
    l'ancien deck en faisait deux slides à un seul paragraphe chacune
    (constat de la revue IC externe sur l'ancienne slide 7) ; le contenu
    qualitatif des deux tient largement sur une slide correctement condensée."""
    slide = _blank_slide(prs)
    _header(slide, "IC Memo — Sections III & V", action_title)

    col_w = (CONTENT_W - 0.4) / 2
    _text(slide, MARGIN_L, CONTENT_TOP, col_w, 0.28, "MARKET & CONSOLIDATION THESIS", size=10, color=COLOR_MUTED, bold=True)
    # Section III est rédigée en prose libre dans le mémo (contrairement à la
    # Section V, une vraie liste à puces) — `_condensed_fragments_from_section`
    # découpe cette prose en fragments courts pour éviter la même "slide à un
    # seul paragraphe" que celle déjà corrigée sur l'Executive Summary.
    market_lines = _adaptive_bullets_from_section(deal.ic_memo, MEMO_SECTIONS[2], limit=6, max_fragments=6)
    _multiline(slide, MARGIN_L, CONTENT_TOP + 0.32, col_w, 4.4, market_lines, size=12.5)

    thesis_left = MARGIN_L + col_w + 0.4
    _text(slide, thesis_left, CONTENT_TOP, col_w, 0.28, "INVESTMENT THESIS", size=10, color=COLOR_MUTED, bold=True)
    thesis_lines = _adaptive_bullets_from_section(deal.ic_memo, MEMO_SECTIONS[4], limit=6, max_fragments=6)
    _multiline(slide, thesis_left, CONTENT_TOP + 0.32, col_w, 4.4, thesis_lines, size=12.5)

    _hairline(slide, MARGIN_L + col_w + 0.2, CONTENT_TOP, 0, color=COLOR_LINE, weight=0.75)
    divider = slide.shapes.add_connector(1, Inches(MARGIN_L + col_w + 0.2), Inches(CONTENT_TOP),
                                          Inches(MARGIN_L + col_w + 0.2), Inches(CONTENT_TOP + 4.6))
    divider.line.color.rgb = COLOR_LINE
    divider.line.width = Pt(0.75)

    comps = ic_context.get("comps_summary")
    if comps:
        _source_note(
            slide, MARGIN_L, FOOTER_TOP - 0.35, CONTENT_W,
            f"Grounded in real listed comparables: {comps['comp_set_name']} (n={comps['sample_size']}), "
            f"median EV/EBITDA {comps['median_ev_ebitda']:.1f}x — see Comparables." if comps.get("median_ev_ebitda") else
            f"Grounded in real listed comparables: {comps['comp_set_name']} (n={comps['sample_size']}).",
        )
    _footer(slide, page, total)


def _slide_comparables(prs: Presentation, comps_table, reference_scenario, ic_context, action_title, page, total) -> None:
    slide = _blank_slide(prs)
    _header(slide, "Market Benchmark", action_title,
            comps_table.comp_set_name if comps_table else "CompSet not available")

    if not comps_table or not comps_table.rows:
        _text(slide, MARGIN_L, CONTENT_TOP + 0.2, 11, 1, "No comparable available for this sector.",
              size=14, color=COLOR_MUTED)
        _footer(slide, page, total)
        return

    rows = [r for r in comps_table.rows if r.ev_ebitda is not None]
    top = sorted(rows, key=lambda r: r.ev_ebitda)[:10]

    comps_summary = (ic_context or {}).get("comps_summary") or {}
    adjacent_tickers = {f["ticker"] for f in (comps_summary.get("adjacent_sector_flags") or [])}
    # Tâche "P3" (Partie D) — `reference_lbo_scenario` (ic_context) n'expose
    # que `entry_multiple_qualified` (texte formaté) ; le multiple NUMÉRIQUE
    # retenu vient du result_json brut du scénario, comme sur la slide
    # Valuation Bridge (avant ce correctif, `.get("entry_multiple")` sur le
    # mauvais dict renvoyait toujours None — la barre cible n'apparaissait
    # jamais sur le football field, constaté à la vérification visuelle).
    entry_mult = (reference_scenario.result_json or {}).get("entry_multiple") if reference_scenario else None

    # Tâche "P3" (Partie B) — vue "football-field" (barres horizontales,
    # ordonnées, cible mise en évidence) au lieu d'un tableau brut. Chaque
    # comparable n'a qu'UNE valeur (pas de fourchette par méthode — donnée
    # non disponible), la cible est ajoutée comme sa propre barre distincte
    # pour la positionner dans la distribution réelle des multiples cotés.
    chart_rows = [(r.ticker, r.ev_ebitda, False) for r in top]
    if entry_mult:
        chart_rows.append(("TARGET (retained)", entry_mult, True))
    chart_rows.sort(key=lambda t: t[1])

    # Tâche "P3" (Partie D) — hauteur de graphique plafonnée plus bas et
    # notes de bas de page resserrées/fusionnées : avec le plein jeu de
    # notes (écart de taille + 2 comparables adjacents), l'ancien calcul
    # débordait sur le pied de page (constaté à la vérification visuelle).
    chart_h = min(3.7, 0.4 * len(chart_rows) + 0.5)
    _add_range_bar_chart(slide, MARGIN_L, CONTENT_TOP, CONTENT_W - 2.0, chart_h, chart_rows)

    legend_x = MARGIN_L + CONTENT_W - 1.85
    _text(slide, legend_x, CONTENT_TOP, 1.85, 0.6,
          "Listed peer\nEV/EBITDA", size=9, color=COLOR_MUTED)
    _rect(slide, legend_x, CONTENT_TOP + 0.55, 0.22, 0.16, fill=COLOR_ACCENT)
    _text(slide, legend_x + 0.3, CONTENT_TOP + 0.48, 1.5, 0.3, "This target", size=9, color=COLOR_INK)

    note_top = CONTENT_TOP + chart_h + 0.15
    median_val = comps_table.stats.median.get("ev_ebitda") if comps_table.stats else None
    if median_val:
        smallest_median = comps_summary.get("smallest_subset_median_ev_ebitda")
        label = f"CompSet median EV/EBITDA: {median_val:.1f}x  (n={len(rows)}, market anchor — not directly applicable)"
        if smallest_median is not None:
            label += f"   |   smallest-subset median: {smallest_median:.1f}x (n=2)"
        _text(slide, MARGIN_L, note_top, CONTENT_W, 0.3, label, size=11.5, color=COLOR_ACCENT, bold=True)
        note_top += 0.32

    size_gap_note = comps_summary.get("size_gap_note")
    if size_gap_note:
        _source_note(slide, MARGIN_L, note_top, CONTENT_W, size_gap_note)
        note_top += 0.46

    adjacent_flags = comps_summary.get("adjacent_sector_flags") or []
    if adjacent_flags:
        combined = "  |  ".join(f"{flag['ticker']}: {flag['note']}" for flag in adjacent_flags[:2])
        _source_note(slide, MARGIN_L, note_top, CONTENT_W, combined)

    _footer(slide, page, total)


def _slide_valuation_bridge(prs: Presentation, deal, comps_table, reference_scenario, ic_context, action_title, page, total) -> None:
    """Tâche "P3" (Partie B) — slide dédiée, demandée explicitement par la
    revue IC externe : le pont médiane comps → décote → multiple retenu,
    en graphique natif (waterfall) plutôt qu'une ligne de texte."""
    slide = _blank_slide(prs)
    _header(slide, "IC Memo — Section VI", action_title, "How the entry multiple was derived")

    if reference_scenario is None:
        _text(slide, MARGIN_L, CONTENT_TOP + 0.2, 11, 1, "No LBO scenario has been modelled for this deal.",
              size=14, color=COLOR_MUTED)
        _footer(slide, page, total)
        return

    r = reference_scenario.result_json or {}
    calib = r.get("calibration") or {}
    entry_mult = r.get("entry_multiple")

    content_top = CONTENT_TOP
    sizing = (ic_context.get("reference_lbo_scenario") or {}).get("sizing_guidance") or {}
    if sizing.get("is_indicative") and sizing.get("note"):
        _text(slide, MARGIN_L, content_top, CONTENT_W, 0.75, sizing["note"],
              size=10.5, color=COLOR_BAD, bold=True)
        content_top += 0.85

    # Tâche "P3" (Partie D) — le bloc KPI est ancré à une position FIXE près
    # du bas de la slide plutôt que calculé en cascade depuis `content_top` :
    # avec une note de taille ET un graphique tous deux présents (petite
    # capitalisation dans un secteur calibré), l'ancien calcul en cascade
    # réduisait la hauteur des tuiles à quelques centièmes de pouce — le
    # texte de valeur débordait alors sur la tuile voisine (constaté à la
    # vérification visuelle : "16." débordant sur "ENTRY EV"). Hauteur fixe,
    # jamais négative ni proche de zéro.
    tile_h = 1.2
    kpi_top = FOOTER_TOP - 0.25 - tile_h - 0.45

    if calib.get("applicable") and calib.get("sufficient"):
        median_mult = calib.get("median_ev_ebitda")
        discount = calib.get("size_illiquidity_discount", 0) * 100
        derived = calib.get("derived_entry_multiple")
        # Seule UNE décote combinée (taille + illiquidité) est réellement
        # calculée par sector_calibration.py — jamais deux décotes séparées
        # qui n'existent pas dans les données ; le libellé de l'étape reflète
        # exactement le champ réel (`discount_label`), pas une invention.
        discount_label = calib.get("discount_label") or "Size & Illiquidity Discount"
        chart_h = min(2.4, kpi_top - 0.45 - content_top)
        _add_bridge_chart(
            slide, MARGIN_L + 0.8, content_top + 0.1, CONTENT_W - 1.6, chart_h,
            [
                ("Comps Median", median_mult, "total"),
                (discount_label, median_mult - derived, "decrease"),
                ("Retained Multiple", derived, "total"),
            ],
            number_format='0.0"x"',
        )
        _source_note(slide, MARGIN_L, content_top + chart_h + 0.2, CONTENT_W,
                     f"CompSet: {calib.get('comp_set_name', 'N/A')} — n={calib.get('sample_size', 'N/A')} "
                     f"comparables with usable real EBITDA. Discount: {discount:.0f}%.")
    else:
        _text(slide, MARGIN_L, content_top + 0.05, CONTENT_W, 0.8,
              calib.get("fallback_reason", "Calibration not applicable — reason not specified."),
              size=13, color=COLOR_MUTED, italic=True)

    fin_prov = r.get("financial_provenance") or {}
    margin_prov = fin_prov.get("ebitda_margin") or {}
    margin_label = {
        "MARKET": "market (comparables)", "DOCUMENT": "documented figure",
        "ESTIMATE": "sector estimate", "REGISTRY": "official registry",
        "MANUAL": "manual entry",
    }.get(margin_prov.get("provenance"), "origin unknown")
    ebitda_margin_val = r.get("ebitda_margin")
    if ebitda_margin_val is None:
        ebitda_margin_val = calib.get("median_ebitda_margin") if calib.get("applicable") and calib.get("sufficient") else None
    if ebitda_margin_val is None:
        from api.services.ma_engine.valuation_engine import LBO_PROFILES, resolve_profile_key
        assumptions = reference_scenario.assumptions_json or {}
        sector_key = resolve_profile_key(assumptions.get("sector_or_naf", ""))
        ebitda_margin_val = LBO_PROFILES[sector_key].ebitda_margin

    kpi_w = (CONTENT_W - 0.6) / 4
    _kpi(slide, MARGIN_L, kpi_top, kpi_w, "Entry Multiple Retained",
         f"{entry_mult:.1f}x" if entry_mult else "N/A", value_color=COLOR_ACCENT, height=tile_h, value_size=22)
    _kpi(slide, MARGIN_L + (kpi_w + 0.2), kpi_top, kpi_w, "EBITDA Margin (assumption)",
         _pct(ebitda_margin_val) if ebitda_margin_val is not None else "N/A",
         value_color=COLOR_INK, sub=margin_label, height=tile_h, value_size=22)
    _kpi(slide, MARGIN_L + 2 * (kpi_w + 0.2), kpi_top, kpi_w, "Entry EV", _money(r.get("entry_ev")),
         value_color=COLOR_INK, height=tile_h, value_size=22)
    _kpi(slide, MARGIN_L + 3 * (kpi_w + 0.2), kpi_top, kpi_w, "Entry Equity", _money(r.get("entry_equity")),
         value_color=COLOR_INK, height=tile_h, value_size=22)

    # Tâche "P3" (Partie C) — la slide ne recopie plus les paragraphes de
    # réconciliation en entier (trop longs pour une slide déjà chargée en
    # graphique/KPI — le débordement l'écraserait) : un renvoi honnête et
    # court vers le mémo, qui porte le texte complet. Jamais un chiffre
    # dissimulé, seulement condensé (voir aussi la note en Section IV/VI du
    # mémo, où le texte complet est reproduit intégralement).
    scenario_ctx = ic_context.get("reference_lbo_scenario") or {}
    if scenario_ctx.get("ebitda_reconciliation_note") or scenario_ctx.get("valuation_reconciliation_note"):
        _source_note(
            slide, MARGIN_L, kpi_top + tile_h + 0.15, CONTENT_W,
            "Note: this scenario's own entry EBITDA/EV differ from the deal-level figures in Financial "
            "Analysis — both are real, on different bases; see the full IC memo (Sections IV & VI) for "
            "the reconciliation.",
        )
    _footer(slide, page, total)


def _slide_financial_analysis(prs: Presentation, deal, ic_context, reference_scenario, action_title, page, total) -> None:
    slide = _blank_slide(prs)
    _header(slide, "IC Memo — Section IV", action_title,
            "Each figure carries its own source, cited alongside it.")

    fin = ic_context["financials"]
    kpi_w = (CONTENT_W - 0.4) / 3
    tile_h = 1.35
    row1 = CONTENT_TOP
    _kpi(slide, MARGIN_L, row1, kpi_w, "Revenue", _money(deal.target_revenue), value_color=COLOR_ACCENT, height=tile_h, value_size=22)
    _kpi(slide, MARGIN_L + kpi_w + 0.2, row1, kpi_w, "EBITDA", _money(deal.target_ebitda), value_color=COLOR_ACCENT, height=tile_h, value_size=22)
    _kpi(slide, MARGIN_L + 2 * (kpi_w + 0.2), row1, kpi_w, "EBITDA Margin",
         _pct(fin["ebitda_margin_pct"], already_pct=True) if fin["ebitda_margin_pct"] is not None else "N/A",
         value_color=COLOR_INK, height=tile_h, value_size=22)

    # Tâche "P3" (Partie B) — barres CA/EBITDA par exercice : le seul
    # historique par "exercice" disponible dans ce produit est la
    # projection annuelle du scénario LBO de référence (même moteur, mêmes
    # chiffres que le mémo/Excel) — jamais une nouvelle série inventée.
    chart_top = row1 + tile_h + 0.3
    projs = (reference_scenario.result_json or {}).get("projections") if reference_scenario else None
    if projs:
        categories = [f"Year {p['year']}" for p in projs]
        _add_column_chart(
            slide, MARGIN_L, chart_top, CONTENT_W, 3.15, categories,
            [("Revenue", [p["revenue"] for p in projs]), ("EBITDA", [p["ebitda"] for p in projs])],
            colors=[COLOR_ACCENT, COLOR_MUTED], number_format='#,##0', has_legend=True,
        )
        _source_note(slide, MARGIN_L, chart_top + 3.2, CONTENT_W,
                     "Projected revenue/EBITDA by year, from the reference LBO scenario — same figures as Section VI/VII.")
    else:
        _multiline(slide, MARGIN_L, chart_top, CONTENT_W, 0.9,
                   ["Quality of Earnings — " + fin["quality_of_earnings"],
                    "Working Capital — " + fin["working_capital"],
                    "Capex Detail — " + fin["capex_detail"]], size=10.5, color=COLOR_MUTED)
    _footer(slide, page, total)


def _slide_capital_structure(prs: Presentation, reference_scenario, action_title, page, total) -> None:
    slide = _blank_slide(prs)
    _header(slide, "IC Memo — Section VI (cont'd)", action_title,
            f"Scenario '{reference_scenario.label}'" if reference_scenario else None)

    if reference_scenario is None:
        _text(slide, MARGIN_L, CONTENT_TOP + 0.2, 11, 1, "No LBO scenario has been modelled for this deal.",
              size=14, color=COLOR_MUTED)
        _footer(slide, page, total)
        return

    r = reference_scenario.result_json or {}
    tranches = r.get("debt_tranches_detail") or []

    # Tâche "P3" (Partie B) — graphique natif Sources & Uses (barres
    # empilées Uses vs Sources) au lieu du rectangle dessiné à la main —
    # mêmes montants que le mémo/Excel (entry_ev, fees, min cash, dette,
    # equity), aucun nouveau calcul.
    # Tâche "P3" (Partie D) — labels désactivés sur les tranches de frais
    # (petites en valeur, leur label chevauchait celui du segment voisin sur
    # une barre empilée serrée — constaté à la vérification visuelle) ; les
    # montants exacts restent disponibles dans le mémo/Excel. Les segments
    # dominants (EV, dette, equity) gardent leur label.
    uses_series = [
        ("Enterprise Value", [r.get("entry_ev") or 0, 0]),
        ("Transaction Fees", [r.get("entry_transaction_fees") or 0, 0], False),
        ("Financing Fees", [r.get("entry_financing_fees") or 0, 0], False),
        ("Minimum Cash", [r.get("entry_min_cash") or 0, 0], False),
    ]
    if tranches:
        source_series = [
            (t.get("name", f"Tranche {i + 1}"), [0, t.get("amount") or 0])
            for i, t in enumerate(tranches)
        ]
    else:
        source_series = [("Senior Debt", [0, r.get("entry_debt") or 0])]
    source_series.append(("Sponsor Equity", [0, r.get("entry_equity") or 0]))
    all_series = uses_series + source_series
    colors = [COLOR_ACCENT, RGBColor(0x8A, 0x8F, 0x99), RGBColor(0xB8, 0xBC, 0xC4), RGBColor(0xD6, 0xD9, 0xDE)] + \
        [RGBColor(0x4A, 0x7A, 0xA8), RGBColor(0x8F, 0xB0, 0xCC), RGBColor(0xC3, 0xD6, 0xE8)][:max(0, len(source_series) - 1)] + [COLOR_INK]

    chart_h = 3.9
    _add_column_chart(
        slide, MARGIN_L + 1.5, CONTENT_TOP, CONTENT_W - 3.0, chart_h, ["Uses", "Sources"],
        all_series, colors=colors, stacked=True, number_format='#,##0', has_legend=True,
    )

    kpi_top = CONTENT_TOP + chart_h + 0.3
    kpi_w = (CONTENT_W - 0.6) / 4
    tile_h = min(1.5, FOOTER_TOP - 0.25 - kpi_top)
    _kpi(slide, MARGIN_L, kpi_top, kpi_w, "Total Uses", _money(r.get("entry_uses_total")),
         value_color=COLOR_INK, height=tile_h, value_size=20)
    _kpi(slide, MARGIN_L + (kpi_w + 0.2), kpi_top, kpi_w, "Leverage at Entry",
         f"{r.get('leverage_entry', 0):.1f}x" if r.get("leverage_entry") is not None else "N/A",
         value_color=COLOR_INK, height=tile_h, value_size=20)
    _kpi(slide, MARGIN_L + 2 * (kpi_w + 0.2), kpi_top, kpi_w, "Sponsor Equity", _money(r.get("entry_equity")),
         value_color=COLOR_ACCENT, height=tile_h, value_size=20)
    equity_pct = (r.get("entry_equity") or 0) / (r.get("entry_uses_total") or 1) * 100
    _kpi(slide, MARGIN_L + 3 * (kpi_w + 0.2), kpi_top, kpi_w, "Equity % of Uses",
         f"{equity_pct:.0f}%", value_color=COLOR_INK, height=tile_h, value_size=20)
    _footer(slide, page, total)


def _slide_deleveraging(prs: Presentation, reference_scenario, action_title, page, total) -> None:
    """Tâche "P3" (Partie B) — NOUVELLE slide : courbe de désendettement
    (levier net Dette nette/EBITDA, et levier brut) sur l'horizon de
    détention. Données lues directement des projections annuelles du
    scénario (net_leverage déjà calculé par le moteur P1, gross leverage =
    debt_eoy/ebitda — une division de deux champs déjà présents, pas une
    nouvelle formule)."""
    slide = _blank_slide(prs)
    _header(slide, "IC Memo — Section VI (cont'd)", action_title, "Net Debt / EBITDA over the holding period")

    projs = (reference_scenario.result_json or {}).get("projections") if reference_scenario else None
    if not projs:
        _text(slide, MARGIN_L, CONTENT_TOP + 0.2, 11, 1, "No LBO scenario has been modelled for this deal.",
              size=14, color=COLOR_MUTED)
        _footer(slide, page, total)
        return

    categories = [f"Year {p['year']}" for p in projs]
    net_leverage = [p.get("net_leverage") for p in projs]
    gross_leverage = [
        (p["debt_eoy"] / p["ebitda"]) if p.get("ebitda") else None
        for p in projs
    ]
    # LINE_MARKERS requires numeric values — None would break the chart XML;
    # Year 0 always has both fields populated in this engine, so this is
    # defensive only.
    if all(v is not None for v in net_leverage) and all(v is not None for v in gross_leverage):
        _add_line_chart(
            slide, MARGIN_L, CONTENT_TOP + 0.1, CONTENT_W, 3.9, categories,
            [("Net Debt / EBITDA", net_leverage), ("Gross Debt / EBITDA", gross_leverage, False)],
            colors=[COLOR_ACCENT, COLOR_MUTED],
        )
    entry_lev = projs[0].get("net_leverage")
    exit_lev = projs[-1].get("net_leverage")
    kpi_top = CONTENT_TOP + 4.2
    kpi_w = (CONTENT_W - 0.4) / 3
    _kpi(slide, MARGIN_L, kpi_top, kpi_w, "Entry Net Leverage", f"{entry_lev:.1f}x" if entry_lev is not None else "N/A",
         value_color=COLOR_INK, height=1.1, value_size=20)
    _kpi(slide, MARGIN_L + kpi_w + 0.2, kpi_top, kpi_w, "Exit Net Leverage", f"{exit_lev:.1f}x" if exit_lev is not None else "N/A",
         value_color=COLOR_GOOD, height=1.1, value_size=20)
    delta = (entry_lev - exit_lev) if entry_lev is not None and exit_lev is not None else None
    _kpi(slide, MARGIN_L + 2 * (kpi_w + 0.2), kpi_top, kpi_w, "Turns Reduced",
         f"{delta:.1f}x" if delta is not None else "N/A", value_color=COLOR_GOOD, height=1.1, value_size=20)
    _footer(slide, page, total)


def _slide_returns_value_creation(prs: Presentation, deal, ic_context, reference_scenario, action_title, page, total) -> None:
    slide = _blank_slide(prs)
    _header(slide, "IC Memo — Section VII", action_title)

    scenario = ic_context.get("reference_lbo_scenario")
    if not scenario:
        _text(slide, MARGIN_L, CONTENT_TOP + 0.2, 11, 1, "No LBO scenario has been modelled for this deal.",
              size=14, color=COLOR_MUTED)
        _footer(slide, page, total)
        return

    r = reference_scenario.result_json or {}
    downside = ic_context.get("downside_scenario")
    half_w = (CONTENT_W - 0.5) / 2

    _text(slide, MARGIN_L, CONTENT_TOP, half_w, 0.28, "BASE vs DOWNSIDE", size=10, color=COLOR_MUTED, bold=True)
    if downside:
        categories = ["Base Case", downside.get("label", "Downside")]
        irr_vals = [scenario.get("irr") or 0, downside.get("irr") or 0]
        moic_vals = [scenario.get("moic") or 0, downside.get("moic") or 0]
        _add_column_chart(
            slide, MARGIN_L, CONTENT_TOP + 0.32, half_w, 2.0, categories,
            [("IRR", irr_vals)], colors=[COLOR_ACCENT], number_format='0.0%', value_axis_min=0,
        )
        _add_column_chart(
            slide, MARGIN_L, CONTENT_TOP + 2.5, half_w, 2.0, categories,
            [("MOIC", moic_vals)], colors=[COLOR_MUTED], number_format='0.00"x"', value_axis_min=0,
        )
    else:
        _text(slide, MARGIN_L, CONTENT_TOP + 0.4, half_w, 0.5,
              "No downside scenario has been modelled for this deal.", size=11, color=COLOR_MUTED, italic=True)

    # Tâche "P3" (Partie B) — pont de création de valeur : Entrée equity →
    # croissance EBITDA → expansion de multiple → désendettement → frais →
    # Sortie equity. Décomposition algébrique EXACTE des chiffres déjà
    # produits par le moteur (voir docstring _add_bridge_chart) — jamais une
    # nouvelle hypothèse : entry/exit EBITDA, entry/exit multiple, entry/exit
    # dette nette et frais sont tous déjà dans result_json.
    bridge_left = MARGIN_L + half_w + 0.5
    _text(slide, bridge_left, CONTENT_TOP, half_w, 0.28, "VALUE CREATION BRIDGE (base case)", size=10, color=COLOR_MUTED, bold=True)
    entry_ebitda, exit_ebitda = r.get("entry_ebitda"), r.get("exit_ebitda")
    entry_mult, exit_mult = r.get("entry_multiple"), r.get("exit_multiple")
    entry_equity, exit_equity = r.get("entry_equity"), r.get("exit_equity")
    entry_debt, entry_min_cash = r.get("entry_debt"), r.get("entry_min_cash") or 0
    exit_net_debt = r.get("exit_net_debt")
    fees = (r.get("entry_transaction_fees") or 0) + (r.get("entry_financing_fees") or 0)
    if None not in (entry_ebitda, exit_ebitda, entry_mult, exit_mult, entry_equity, exit_equity, entry_debt, exit_net_debt):
        entry_net_debt = entry_debt - entry_min_cash
        ebitda_growth = (exit_ebitda - entry_ebitda) * entry_mult
        multiple_expansion = exit_ebitda * (exit_mult - entry_mult)
        deleveraging = entry_net_debt - exit_net_debt

        def _step(label, delta):
            return (label, abs(delta), "increase" if delta >= 0 else "decrease")

        steps = [
            ("Entry Equity", entry_equity, "total"),
            _step("EBITDA\nGrowth", ebitda_growth),
            _step("Multiple\nExpansion", multiple_expansion),
            _step("Deleveraging", deleveraging),
            _step("Fees", -fees),
            ("Exit Equity", exit_equity, "total"),
        ]
        # Tâche "P3" (Partie D) — le label du dernier bâton (Exit Equity, le
        # plus grand montant) débordait du cadre du graphique en écriture
        # intégrale (constaté à la vérification visuelle) : format compact
        # €M/€K adapté à l'ordre de grandeur réel du deal, jamais un montant
        # arrondi différemment de celui affiché ailleurs (juste une échelle
        # d'affichage — la valeur porteuse reste exacte).
        max_abs = max(abs(v) for _, v, _ in steps)
        bridge_number_format = '#,##0.0,,"M"' if max_abs >= 10_000_000 else '#,##0,"K"'
        _add_bridge_chart(slide, bridge_left, CONTENT_TOP + 0.32, half_w, 4.2, steps, number_format=bridge_number_format)
    else:
        _text(slide, bridge_left, CONTENT_TOP + 0.4, half_w, 0.5,
              "Value creation bridge not available — incomplete scenario data.", size=11, color=COLOR_MUTED, italic=True)
    _footer(slide, page, total)


def _slide_sensitivity(prs: Presentation, deal, ic_context, action_title, page, total) -> None:
    """Grille de sensibilité IRR/MOIC — reléguée à sa propre slide (Tâche
    "P3") pour ne pas surcharger la slide Returns & Value Creation."""
    slide = _blank_slide(prs)
    _header(slide, "IC Memo — Section VII (cont'd)", action_title)
    sensitivity = ic_context.get("sensitivity")
    if not sensitivity:
        _text(slide, MARGIN_L, CONTENT_TOP + 0.2, 11, 1,
              "Sensitivity analysis not available for this scenario.", size=14, color=COLOR_MUTED, italic=True)
        _footer(slide, page, total)
        return

    headers = ["Leverage \\ Exit"] + [f"{ex:.1f}x" for ex in sensitivity["exit_axis"]]
    data = []
    for row in sensitivity["grid"]:
        cells = [f"{row['leverage']:.1f}x"]
        for cell in row["cells"]:
            irr_txt = _pct(cell["irr"]) if cell["irr"] is not None else "N/A"
            moic_txt = f"{cell['moic']:.2f}x" if cell["moic"] is not None else "N/A"
            cells.append(f"{irr_txt} / {moic_txt}")
        data.append(cells)
    table_h = 0.42 * (len(data) + 1)
    _styled_table(
        slide, len(data) + 1, len(headers), MARGIN_L, CONTENT_TOP + 0.3, CONTENT_W, table_h,
        headers=headers, data=data,
        col_align=["left"] + ["right"] * (len(headers) - 1), header_size=10.5, body_size=10.5,
    )
    _source_note(slide, MARGIN_L, CONTENT_TOP + 0.3 + table_h + 0.15, CONTENT_W, sensitivity["method"])
    _footer(slide, page, total)


def _slide_risk_factors(prs: Presentation, deal, ic_context, action_title, page, total) -> None:
    slide = _blank_slide(prs)
    _header(slide, "IC Memo — Section VIII", action_title)
    # Tâche "P3" (Partie D) — `_marked_bullets_from_section` : le ⚠
    # n'apparaît que sur les puces réelles, jamais sur une ligne d'intro en
    # prose libre (constat de la revue IC externe).
    lines = _marked_bullets_from_section(deal.ic_memo, MEMO_SECTIONS[7], marker="⚠", limit=9)
    _multiline(slide, MARGIN_L, CONTENT_TOP + 0.1, CONTENT_W, 3.9, lines, size=13)

    self_check = ic_context.get("self_check") or {}
    failed = [c for c in (self_check.get("checks") or []) if not c["passed"]]
    if failed:
        note_top = CONTENT_TOP + 4.1
        _text(slide, MARGIN_L, note_top, CONTENT_W, 0.28, "KNOWN DATA / MODEL LIMITATIONS", size=10, color=COLOR_BAD, bold=True)
        _multiline(slide, MARGIN_L, note_top + 0.3, CONTENT_W, 1.4,
                   [f"⚠ {c['name']}: {c['detail']}" for c in failed], size=10, color=COLOR_MUTED)
    _footer(slide, page, total)


def _slide_recommendation(prs: Presentation, deal, ic_context, action_title, page, total) -> None:
    """Tâche "P3" (Partie C) — slide de recommandation ORIENTÉE ACTION :
    l'ancien deck se contentait de coller la phrase du mémo sur une page
    vide (constat de la revue IC externe, ancienne slide 12). Reprend les
    chiffres clés (déjà affichés ailleurs dans le même deck) à côté de
    l'appel et des conditions, pour qu'un comité voie tout sur une slide."""
    slide = _blank_slide(prs)
    _header(slide, "IC Memo — Section IX", action_title)

    scenario = ic_context.get("reference_lbo_scenario") or {}
    downside = ic_context.get("downside_scenario")
    kpi_w = (CONTENT_W - 0.6) / 4
    tile_h = 1.2
    irr, moic = scenario.get("irr"), scenario.get("moic")
    irr_color = COLOR_GOOD if irr and irr >= 0.20 else COLOR_BAD if irr is not None else COLOR_MUTED
    moic_color = COLOR_GOOD if moic and moic >= 2.5 else COLOR_BAD if moic is not None else COLOR_MUTED
    _kpi(slide, MARGIN_L, CONTENT_TOP, kpi_w, "Base IRR", _pct(irr) if irr is not None else "N/A", value_color=irr_color, height=tile_h, value_size=22)
    _kpi(slide, MARGIN_L + (kpi_w + 0.2), CONTENT_TOP, kpi_w, "Base MOIC", f"{moic:.2f}x" if moic is not None else "N/A", value_color=moic_color, height=tile_h, value_size=22)
    _kpi(slide, MARGIN_L + 2 * (kpi_w + 0.2), CONTENT_TOP, kpi_w, "Downside IRR",
         _pct(downside.get("irr")) if downside and downside.get("irr") is not None else "N/A",
         value_color=COLOR_MUTED, height=tile_h, value_size=22)
    _kpi(slide, MARGIN_L + 3 * (kpi_w + 0.2), CONTENT_TOP, kpi_w, "Downside MOIC",
         f"{downside['moic']:.2f}x" if downside and downside.get("moic") is not None else "N/A",
         value_color=COLOR_MUTED, height=tile_h, value_size=22)

    lines = _bullets_from_section(deal.ic_memo, MEMO_SECTIONS[8], limit=10)
    _multiline(slide, MARGIN_L, CONTENT_TOP + tile_h + 0.35, CONTENT_W, 2.6, lines, size=15)

    next_steps_top = CONTENT_TOP + tile_h + 0.35 + 2.2
    _text(slide, MARGIN_L, next_steps_top, CONTENT_W, 0.28, "NEXT STEPS", size=10, color=COLOR_MUTED, bold=True)
    rec_guidance = ic_context.get("recommendation_guidance") or {}
    next_steps = (
        ["Confirm EBITDA via audited financials before proceeding further."]
        if not rec_guidance.get("ebitda_is_real")
        else ["Proceed to detailed due diligence and financing arrangement."]
    )
    next_steps.append("Revisit sensitivity and downside case as diligence findings come in.")
    _multiline(slide, MARGIN_L, next_steps_top + 0.32, CONTENT_W, 0.9,
               [f"→ {s}" for s in next_steps], size=12, color=COLOR_INK)
    _footer(slide, page, total)


def _slide_appendix_provenance(prs: Presentation, deal, ic_context, page, total) -> None:
    """Tâche "P3" (Partie C) — annexe : provenance des chiffres clés en
    langage naturel (cohérent avec la purge P0 — jamais un chemin de fichier,
    un tag interne ou un message d'erreur français dans ce deck anglais)."""
    slide = _blank_slide(prs)
    _header(slide, "Appendix", "Data Provenance", "Where each key figure comes from")

    prov = deal.financial_provenance if isinstance(deal.financial_provenance, dict) else {}
    scenario_ctx = ic_context.get("reference_lbo_scenario") or {}
    rows = []
    for label, field in [("Revenue", "target_revenue"), ("EBITDA", "target_ebitda"), ("Enterprise Value", "enterprise_value")]:
        entry = prov.get(field) or {}
        rows.append([label, entry.get("provenance", "UNKNOWN"), entry.get("reference") or "—"])
    rows.append(["Entry Multiple", "ESTIMATE" if scenario_ctx else "N/A", scenario_ctx.get("entry_multiple_qualified", "—") if scenario_ctx else "—"])

    # Tâche "P3" (Partie D) — plus de troncature en plein mot : la colonne
    # "Reference" (souvent 200+ caractères) reçoit le texte complet, dans une
    # colonne large dédiée qui limite l'enroulement à ~3 lignes ; row_h=0.85
    # laisse ces 3 lignes s'afficher sans que le tableau ne déborde sur le
    # footer (5 lignes × row_h doit tenir entre CONTENT_TOP et FOOTER_TOP).
    row_h = 0.85
    _styled_table(
        slide, len(rows) + 1, 3, MARGIN_L, CONTENT_TOP, CONTENT_W, row_h * (len(rows) + 1),
        headers=["Figure", "Provenance", "Reference"], data=rows,
        col_align=["left", "left", "left"], body_size=10,
        col_widths=[1.6, 1.5, CONTENT_W - 1.6 - 1.5],
    )
    _source_note(
        slide, MARGIN_L, CONTENT_TOP + row_h * (len(rows) + 1) + 0.25, CONTENT_W,
        "REGISTRY/DOCUMENT/MARKET = real, sourced data. ESTIMATE = modelled from sector assumptions, "
        "not yet confirmed by audited financials. See the IC memo for the full per-figure citation.",
    )
    _footer(slide, page, total)


# ============================================================
# Point d'entrée
# ============================================================

def _build_action_titles(deal, ic_context, reference_scenario, comps_table) -> dict[str, str]:
    """Tâche "P3 : le deck IC devient un vrai deck" (Partie A) — un titre
    d'action énonce le MESSAGE de la slide, jamais l'intitulé de section
    ("V. Investment Thesis"). Chaque titre ci-dessous est calculé depuis les
    chiffres RÉELS déjà présents dans `ic_context`/`result_json` — jamais un
    nombre inventé pour l'occasion ; en l'absence de données, un titre
    neutre (mais toujours orienté message) est utilisé plutôt qu'un titre
    de section brut."""
    scenario = ic_context.get("reference_lbo_scenario") or {}
    downside = ic_context.get("downside_scenario")
    r = (reference_scenario.result_json or {}) if reference_scenario else {}
    calib = r.get("calibration") or {}
    sector = deal.sector or "the sector"
    titles: dict[str, str] = {}

    irr, moic = scenario.get("irr"), scenario.get("moic")
    n = scenario.get("holding_period_years")
    if irr is not None and moic is not None and n:
        titles["executive_summary"] = f"{moic:.1f}x MOIC / {irr * 100:.0f}% IRR projected over {int(n)} years"
    else:
        titles["executive_summary"] = "Deal snapshot: revenue, EBITDA and projected returns"

    titles["company_overview"] = f"{deal.target_name or 'The target'}: a {sector} platform"

    titles["market_thesis"] = f"Consolidation opportunity in a fragmented {sector} market"

    median_ev_ebitda = (ic_context.get("comps_summary") or {}).get("median_ev_ebitda")
    if median_ev_ebitda:
        titles["comparables"] = f"Public peers trade at {median_ev_ebitda:.1f}x EV/EBITDA — a market anchor only"
    else:
        titles["comparables"] = "No listed comparables available for this sector"

    if calib.get("applicable") and calib.get("sufficient"):
        derived = calib.get("derived_entry_multiple")
        discount = (calib.get("size_illiquidity_discount") or 0) * 100
        titles["valuation_bridge"] = f"{calib.get('median_ev_ebitda', 0):.1f}x comps median discounted {discount:.0f}% to {derived:.1f}x retained"
    else:
        titles["valuation_bridge"] = "Entry multiple set by generic sector profile — no real comparables available"

    margin_pct = ic_context["financials"].get("ebitda_margin_pct")
    if margin_pct is not None:
        titles["financial_analysis"] = f"{_money(deal.target_revenue)} revenue at a {margin_pct:.0f}% EBITDA margin"
    else:
        titles["financial_analysis"] = "Revenue and EBITDA — sourced and qualified figures"

    leverage = r.get("leverage_entry")
    uses_total = r.get("entry_uses_total")
    equity = r.get("entry_equity")
    if leverage is not None and uses_total:
        equity_pct = (equity or 0) / uses_total * 100
        titles["capital_structure"] = f"{leverage:.1f}x leverage funds a {_money(uses_total)} deal, {equity_pct:.0f}% equity"
    else:
        titles["capital_structure"] = "Sources & Uses of financing"

    projs = r.get("projections") or []
    if projs and projs[0].get("net_leverage") is not None and projs[-1].get("net_leverage") is not None:
        titles["deleveraging"] = f"Net leverage falls from {projs[0]['net_leverage']:.1f}x to {projs[-1]['net_leverage']:.1f}x by exit"
    else:
        titles["deleveraging"] = "Deleveraging profile over the holding period"

    if moic is not None and downside and downside.get("moic") is not None:
        titles["returns"] = f"{moic:.1f}x base vs {downside['moic']:.2f}x downside case"
    elif moic is not None:
        titles["returns"] = f"{moic:.1f}x projected MOIC — value creation bridge"
    else:
        titles["returns"] = "Returns Analysis"

    titles["sensitivity"] = "Returns hold up across exit and leverage assumptions"

    risk_candidates = ic_context.get("risk_candidates") or []
    if risk_candidates:
        titles["risk_factors"] = f"{len(risk_candidates)} sector-specific risks identified, prioritized by severity"
    else:
        titles["risk_factors"] = "Key Risk Factors"

    rec_guidance = ic_context.get("recommendation_guidance") or {}
    titles["recommendation"] = f"Recommendation: {rec_guidance.get('required_positive_label', 'See narrative below')}"

    return titles


def generate_ic_deck(deal, comps_table=None, reference_scenario=None, ic_context: dict | None = None) -> io.BytesIO:
    """Generates the full IC deck for a promoted deal — a genuine committee
    deck (action titles, native charts, condensed points), not the memo
    paginated (Tâche "P3 : le deck IC devient un vrai deck"). Consumes the
    SAME structured context (`ic_context.build_ic_context`) as the Word
    memo so the two documents can never disagree on a figure.

    Args:
        deal: `Deal` instance.
        comps_table: `CompsTableResponse | None` (Comps Engine, calibrated sector).
        reference_scenario: `LBOScenario | None` (D23) — reference scenario.
        ic_context: structured context from `ic_context.build_ic_context` —
            pass the SAME object used for the Word memo so the two documents
            can never disagree on a figure.

    Returns:
        io.BytesIO ready to stream (.pptx).
    """
    logger.info("📊 Generating IC deck for deal #{} ({})…", deal.id, deal.target_name)

    if ic_context is None:
        ic_context = build_ic_context(deal, reference_scenario, comps_table)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    titles = _build_action_titles(deal, ic_context, reference_scenario, comps_table)

    agenda_items = [
        "Executive Summary",
        "Company Overview",
        "Market & Investment Thesis",
        "Comparables & Valuation Bridge",
        "Financial Analysis",
        "Capital Structure & Deleveraging",
        "Returns Analysis & Value Creation",
        "Risk Factors",
        "Recommendation",
    ]

    total = 15
    _slide_cover(prs, deal)
    _slide_agenda(prs, agenda_items, 2, total)
    _slide_executive_summary(prs, deal, ic_context, titles["executive_summary"], 3, total)
    _slide_company_overview(prs, deal, ic_context, titles["company_overview"], 4, total)
    _slide_industry_market_thesis(prs, deal, ic_context, titles["market_thesis"], 5, total)
    _slide_comparables(prs, comps_table, reference_scenario, ic_context, titles["comparables"], 6, total)
    _slide_valuation_bridge(prs, deal, comps_table, reference_scenario, ic_context, titles["valuation_bridge"], 7, total)
    _slide_financial_analysis(prs, deal, ic_context, reference_scenario, titles["financial_analysis"], 8, total)
    _slide_capital_structure(prs, reference_scenario, titles["capital_structure"], 9, total)
    _slide_deleveraging(prs, reference_scenario, titles["deleveraging"], 10, total)
    _slide_returns_value_creation(prs, deal, ic_context, reference_scenario, titles["returns"], 11, total)
    _slide_sensitivity(prs, deal, ic_context, titles["sensitivity"], 12, total)
    _slide_risk_factors(prs, deal, ic_context, titles["risk_factors"], 13, total)
    _slide_recommendation(prs, deal, ic_context, titles["recommendation"], 14, total)
    _slide_appendix_provenance(prs, deal, ic_context, 15, total)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    logger.info("  ✅ IC deck generated ({:.0f} KB, {} slides).", buf.getbuffer().nbytes / 1024, len(prs.slides))
    return buf
